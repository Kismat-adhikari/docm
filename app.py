import os
from dotenv import load_dotenv
import json
import requests
from flask import Flask, render_template, request, redirect, url_for, flash, jsonify, session
from werkzeug.utils import secure_filename
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from quizpasa import handle_chat_message, get_welcome_message
from flask import session


import docx2txt
from docx import Document
from pptx import Presentation
import tempfile
import shutil
from PIL import Image, ImageEnhance, ImageFilter
import io
import subprocess
import sys
import fitz  # PyMuPDF - better for image extraction from PDFs
import zipfile
from io import BytesIO
import base64

app = Flask(__name__)
app.config['SECRET_KEY'] = 'your-secret-key-here'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Configure session for temporary storage
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_PERMANENT'] = False
app.config['SESSION_USE_SIGNER'] = True
app.config['SESSION_KEY_PREFIX'] = 'studypasa:'



# Ensure uploads directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'pptx', 'png', 'jpg', 'jpeg', 'tiff', 'bmp'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS



def validate_file_content(file_path, file_ext):
    """Validate file content before processing"""
    try:
        # Check if file is empty
        if os.path.getsize(file_path) == 0:
            return False, 'The file is empty'
            
        # Basic validation based on file type
        if file_ext in ['doc', 'docx']:
            try:
                if file_ext == 'doc':
                    # Use docx2txt for .doc files
                    text = docx2txt.process(file_path)
                    if not text or not text.strip():
                        return False, 'The document contains no readable text'
                else:
                    # Use python-docx for .docx files
                    doc = Document(file_path)
                    has_content = False
                    
                    # Check paragraphs
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            has_content = True
                            break
                    
                    # Check tables if no paragraph content
                    if not has_content:
                        for table in doc.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        has_content = True
                                        break
                                if has_content:
                                    break
                            if has_content:
                                break
                    
                    if not has_content:
                        return False, 'The document contains no readable text'
                        
                return True, ''
                
            except Exception as e:
                return False, f'The document appears to be corrupted or invalid: {str(e)}'
                
        elif file_ext == 'pdf':
            try:
                with pdfplumber.open(file_path) as pdf:
                    if len(pdf.pages) == 0:
                        return False, 'The PDF contains no pages'
                    return True, ''
            except Exception as e:
                return False, f'The PDF appears to be corrupted or password protected: {str(e)}'
                
        elif file_ext == 'pptx':
            try:
                prs = Presentation(file_path)
                if len(prs.slides) == 0:
                    return False, 'The presentation contains no slides'
                # Check if any slide has text
                has_text = False
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            has_text = True
                            break
                    if has_text:
                        break
                if not has_text:
                    return False, 'The presentation contains no readable text'
                return True, ''
            except Exception as e:
                return False, f'The presentation appears to be corrupted or invalid: {str(e)}'
                    
        elif file_ext in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            try:
                with Image.open(file_path) as img:
                    # Check if image is not completely blank or black
                    extrema = img.convert('L').getextrema()
                    if extrema[0] == extrema[1]:  # All pixels are the same value
                        return False, 'The image appears to be blank or contains no content'
                    # Check if image is too small for reliable OCR
                    if img.size[0] < 50 or img.size[1] < 50:
                        return False, 'The image is too small for text extraction'
                    # For images, allow validation to pass so OCR can be applied
                    return True, ''
            except Exception as e:
                return False, f'The image appears to be corrupted or in an invalid format: {str(e)}'
                
    except Exception as e:
        return False, f'Error validating file: {str(e)}'


def setup_tesseract():
    """Enhanced Tesseract setup with better path detection"""
    try:
        # First try the default command
        result = subprocess.run(['tesseract', '--version'], 
                              capture_output=True, text=True, timeout=10)
        if result.returncode == 0:
            print(f"Tesseract found: {result.stdout.split()[1]}")
            return True
    except:
        pass
    
    # Try common installation paths
    possible_paths = [
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        r'C:\Users\{}\AppData\Local\Tesseract-OCR\tesseract.exe'.format(os.getenv('USERNAME', '')),
        '/usr/bin/tesseract',
        '/usr/local/bin/tesseract',
        '/opt/homebrew/bin/tesseract',
        '/usr/local/Cellar/tesseract/*/bin/tesseract'  # Homebrew on macOS
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            try:
                pytesseract.pytesseract.tesseract_cmd = path
                result = subprocess.run([path, '--version'], 
                                      capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    print(f"Tesseract found at: {path}")
                    return True
            except:
                continue
    
    print("Tesseract not found. Please install it:")
    print("Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki")
    print("Linux: sudo apt-get install tesseract-ocr")
    print("macOS: brew install tesseract")
    return False

def preprocess_image_for_ocr(image):
    """Enhance image for better OCR results"""
    try:
        # Convert to grayscale
        if image.mode != 'L':
            image = image.convert('L')
        
        # Enhance contrast
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        
        # Enhance sharpness
        enhancer = ImageEnhance.Sharpness(image)
        image = enhancer.enhance(2.0)
        
        # Apply slight blur to reduce noise
        image = image.filter(ImageFilter.MedianFilter(size=3))
        
        return image
    except Exception as e:
        print(f"Error preprocessing image: {e}")
        return image

def extract_text_from_image(file_path):
    """Enhanced image text extraction with better OCR"""
    try:
        print(f"Extracting text from image: {file_path}")
        
        if not setup_tesseract():
            flash("Tesseract OCR is not installed. Please install it to process images.", 'error')
            return ""
        
        with Image.open(file_path) as img:
            print(f"Original image - Mode: {img.mode}, Size: {img.size}")
            
            # Convert to RGB if necessary
            if img.mode not in ['L', 'RGB']:
                img = img.convert('RGB')
            
            # Preprocess image for better OCR
            processed_img = preprocess_image_for_ocr(img)
            
            # Try multiple OCR configurations
            ocr_configs = [
                '--psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz.,!?;: ',
                '--psm 4 --oem 3',
                '--psm 3 --oem 3',
                '--psm 6 --oem 3',
                '--psm 1 --oem 3',
                '--psm 8 --oem 3',
                '--psm 13 --oem 3'
            ]
            
            best_text = ""
            best_length = 0
            
            for config in ocr_configs:
                try:
                    text = pytesseract.image_to_string(processed_img, lang='eng', config=config)
                    if len(text.strip()) > best_length:
                        best_text = text
                        best_length = len(text.strip())
                        print(f"Better result with config '{config}': {best_length} chars")
                except Exception as e:
                    print(f"OCR config '{config}' failed: {e}")
                    continue
            
            # If still no good result, try with original image
            if best_length < 10:
                try:
                    text = pytesseract.image_to_string(img, lang='eng')
                    if len(text.strip()) > best_length:
                        best_text = text
                        print(f"Better result with original image: {len(text.strip())} chars")
                except Exception as e:
                    print(f"Original image OCR failed: {e}")
            
            print(f"Final OCR result: {len(best_text)} characters")
            return best_text.strip()
    
    except Exception as e:
        print(f"Error in extract_text_from_image: {e}")
        return ""

def extract_images_from_pdf(file_path):
    """Extract images from PDF and perform OCR"""
    text = ""
    
    try:
        doc = fitz.open(file_path)
        print(f"PDF has {len(doc)} pages")
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images()
            
            print(f"Page {page_num + 1} has {len(image_list)} images")
            
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    # Convert to PIL Image
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("ppm")
                        pil_image = Image.open(io.BytesIO(img_data))
                        
                        # Preprocess and OCR
                        processed_img = preprocess_image_for_ocr(pil_image)
                        
                        # Extract text from image
                        img_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                        
                        if img_text.strip():
                            text += f"\n[Image {img_index + 1} from Page {page_num + 1}]:\n{img_text}\n"
                            print(f"Extracted {len(img_text)} chars from image {img_index + 1} on page {page_num + 1}")
                    
                    pix = None
                    
                except Exception as e:
                    print(f"Error processing image {img_index + 1} on page {page_num + 1}: {e}")
                    continue
        
        doc.close()
        
    except Exception as e:
        print(f"Error extracting images from PDF: {e}")
    
    return text

def extract_images_from_docx(file_path):
    """Extract images from DOCX and perform OCR"""
    text = ""
    
    try:
        # DOCX is essentially a ZIP file
        with zipfile.ZipFile(file_path, 'r') as docx_zip:
            # Look for image files in the media folder
            image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
            
            print(f"Found {len(image_files)} images in DOCX")
            
            for img_file in image_files:
                try:
                    # Extract image data
                    img_data = docx_zip.read(img_file)
                    
                    # Convert to PIL Image
                    pil_image = Image.open(io.BytesIO(img_data))
                    
                    # Preprocess and OCR
                    processed_img = preprocess_image_for_ocr(pil_image)
                    
                    # Extract text
                    img_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                    
                    if img_text.strip():
                        text += f"\n[Image from DOCX: {img_file}]:\n{img_text}\n"
                        print(f"Extracted {len(img_text)} chars from {img_file}")
                
                except Exception as e:
                    print(f"Error processing image {img_file}: {e}")
                    continue
    
    except Exception as e:
        print(f"Error extracting images from DOCX: {e}")
    
    return text

def extract_images_from_pptx(file_path):
    """Extract images from PPTX and perform OCR"""
    text = ""
    
    try:
        # PPTX is also a ZIP file
        with zipfile.ZipFile(file_path, 'r') as pptx_zip:
            # Look for image files in the media folder
            image_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/media/')]
            
            print(f"Found {len(image_files)} images in PPTX")
            
            for img_file in image_files:
                try:
                    # Extract image data
                    img_data = pptx_zip.read(img_file)
                    
                    # Convert to PIL Image
                    pil_image = Image.open(io.BytesIO(img_data))
                    
                    # Preprocess and OCR
                    processed_img = preprocess_image_for_ocr(pil_image)
                    
                    # Extract text
                    img_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                    
                    if img_text.strip():
                        text += f"\n[Image from PPTX: {img_file}]:\n{img_text}\n"
                        print(f"Extracted {len(img_text)} chars from {img_file}")
                
                except Exception as e:
                    print(f"Error processing image {img_file}: {e}")
                    continue
    
    except Exception as e:
        print(f"Error extracting images from PPTX: {e}")
    
    return text

def extract_text_from_pdf(file_path):
    """Enhanced PDF text extraction with image OCR"""
    text = ""
    
    try:
        print(f"Processing PDF: {file_path}")
        
        # First, extract regular text with pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                    print(f"Page {page_num + 1}: Extracted {len(page_text)} characters with pdfplumber")
        
        # Then extract text from images in the PDF
        if setup_tesseract():
            image_text = extract_images_from_pdf(file_path)
            if image_text:
                text += "\n=== TEXT FROM IMAGES ===\n" + image_text
        
        # If still no text, try OCR on the entire PDF pages
        if not text.strip() or len(text.strip()) < 50:
            print("Minimal text found, trying full page OCR...")
            
            if setup_tesseract():
                try:
                    images = convert_from_path(file_path, dpi=300)
                    ocr_text = ""
                    
                    for i, image in enumerate(images):
                        processed_img = preprocess_image_for_ocr(image)
                        page_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                        ocr_text += f"\n[Page {i + 1} OCR]:\n{page_text}\n"
                        print(f"OCR Page {i + 1}: Extracted {len(page_text)} characters")
                    
                    if ocr_text.strip():
                        text = ocr_text
                        
                except Exception as e:
                    print(f"Full page OCR failed: {e}")
    
    except Exception as e:
        print(f"Error processing PDF: {e}")
    
    return text


def extract_text_from_docx(file_path):
    """Enhanced DOCX/DOC text extraction with image OCR"""
    text = ""
    
    try:
        file_ext = file_path.rsplit('.', 1)[1].lower()
        print(f"Processing {file_ext.upper()}: {file_path}")
        
        if file_ext == 'doc':
            # Handle .doc files using docx2txt
            text = docx2txt.process(file_path)
            print(f"DOC text extracted: {len(text)} characters")
        elif file_ext == 'docx':
            # Handle .docx files using python-docx (for better table support)
            doc = Document(file_path)
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():  # Only add non-empty paragraphs
                    text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():  # Only add non-empty cells
                            text += cell.text + "\n"
            
            print(f"DOCX text extracted: {len(text)} characters")
        
        # Extract text from images in DOCX (only for .docx files)
        if file_ext == 'docx' and setup_tesseract():
            try:
                image_text = extract_images_from_docx(file_path)
                if image_text:
                    text += "\n=== TEXT FROM IMAGES ===\n" + image_text
            except Exception as e:
                print(f"Error extracting images from DOCX: {e}")
        
        # If no text was extracted, provide a helpful message
        if not text.strip():
            print(f"No text extracted from {file_ext.upper()} file")
            return ""
        
        return text.strip()
        
    except Exception as e:
        print(f"Error extracting text from DOCX/DOC: {e}")
        return ""

    
def extract_text_from_pptx(file_path):
    """Enhanced PPTX text extraction with image OCR"""
    text = ""
    
    try:
        print(f"Processing PPTX: {file_path}")
        
        # Extract regular text
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n=== SLIDE {slide_num + 1} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        
        print(f"PPTX text extracted: {len(text)} characters")
        
        # Extract text from images
        if setup_tesseract():
            image_text = extract_images_from_pptx(file_path)
            if image_text:
                text += "\n=== TEXT FROM IMAGES ===\n" + image_text
    
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    
    return text

def generate_mcqs_with_groq(text, num_questions=10):
    """Generate MCQs using GROQ API"""
    # Load environment variables
    load_dotenv()
    # GROQ API configuration
    GROQ_API_KEY = os.getenv("GROQ_API_KEY", "your-groq-api-key-here")
    GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
    # If no API key is provided, use fallback
    if not GROQ_API_KEY or GROQ_API_KEY == "your-groq-api-key-here":
        print("No GROQ API key provided, using fallback MCQ generation")
        return create_content_based_mcqs(text, num_questions)
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }
    
    # Extract key topics and themes from the text
    text_preview = text[:4000] if len(text) > 4000 else text
    
    prompt = f"""
    You are an expert question generator. Based on the following text content, create {num_questions} high-quality multiple choice questions (MCQs).

    CONTENT TO ANALYZE:
    {text_preview}

    REQUIREMENTS:
    1. Questions must be directly related to the content above
    2. Cover different topics and concepts from the text
    3. Mix of difficulty levels (easy, medium, hard)
    4. Each question has exactly 4 options
    5. Only ONE correct answer per question
    6. Distractors should be plausible but clearly wrong
    7. Focus on key concepts, facts, and important details

    RESPONSE FORMAT (JSON only):
    {{
        "questions": [
            {{
                "question": "What is the main concept discussed regarding [specific topic]?",
                "options": [
                    "Correct answer based on text",
                    "Plausible but incorrect option",
                    "Another plausible but incorrect option",
                    "Third plausible but incorrect option"
                ],
                "correct_answer": 0
            }}
        ]
    }}

    Generate {num_questions} questions now. Return ONLY the JSON response, no additional text.
    """
    
    payload = {
        "model": "llama3-8b-8192",
        "messages": [
            {
                "role": "system",
                "content": "You are an expert educational content creator specializing in generating high-quality multiple choice questions from text content. Always respond with valid JSON format only."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        "temperature": 0.3,
        "max_tokens": 2500
    }
    
    try:
        print(f"Calling GROQ API to generate {num_questions} MCQs...")
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)
        
        if response.status_code == 200:
            result = response.json()
            content = result['choices'][0]['message']['content'].strip()
            
            # Clean up the response
            if content.startswith('```json'):
                content = content.replace('```json', '').replace('```', '').strip()
            
            try:
                mcq_data = json.loads(content)
                if 'questions' in mcq_data and len(mcq_data['questions']) > 0:
                    print(f"Successfully generated {len(mcq_data['questions'])} MCQs")
                    return mcq_data['questions']
                else:
                    print("No questions found in API response")
                    return create_content_based_mcqs(text, num_questions)
            except json.JSONDecodeError as e:
                print(f"JSON parsing error: {e}")
                return create_content_based_mcqs(text, num_questions)
        else:
            print(f"GROQ API error: {response.status_code} - {response.text}")
            return create_content_based_mcqs(text, num_questions)
            
    except Exception as e:
        print(f"Error calling GROQ API: {e}")
        return create_content_based_mcqs(text, num_questions)

def create_content_based_mcqs(text, num_questions=10):
    """Create MCQs based on actual text content"""
    
    words = text.split()
    sentences = [s.strip() for s in text.split('.') if s.strip()]
    paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
    
    questions = []
    
    # Question 1: About document content
    if len(sentences) > 0:
        first_sentence = sentences[0][:100] + "..." if len(sentences[0]) > 100 else sentences[0]
        questions.append({
            "question": f"According to the document, what is mentioned in the opening section?",
            "options": [
                f"Content related to: {first_sentence}",
                "Financial planning strategies",
                "Sports team statistics",
                "Weather forecast data"
            ],
            "correct_answer": 0
        })
    
    # Question 2: About document length
    word_count = len(words)
    questions.append({
        "question": f"Approximately how many words does this document contain?",
        "options": [
            f"Around {word_count} words",
            f"Around {word_count//2} words",
            f"Around {word_count*2} words",
            "Less than 100 words"
        ],
        "correct_answer": 0
    })
    
    # Add more questions following the same pattern...
    # (Include the rest of the original create_content_based_mcqs function here)
    
    return questions[:num_questions]

def generate_flashcards_with_groq(text, num_flashcards=10):
    load_dotenv()
    GROQ_API_KEY = os.getenv("GROQ_API_KEY", "your-groq-api-key-here")
    GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"

    if not GROQ_API_KEY or GROQ_API_KEY == "your-groq-api-key-here":
        print("No GROQ API key provided for flashcards, using fallback generation")
        return create_content_based_flashcards(text, num_flashcards)

    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    text_preview = text[:4000] if len(text) > 4000 else text

    prompt = f"""
    You are an expert flashcard generator. Based on the following text content, create {num_flashcards} high-quality flashcards.

    CONTENT TO ANALYZE:
    {text_preview}

    REQUIREMENTS:
    1. Each flashcard must have a clear 'front' (question/term) and 'back' (answer/definition).
    2. Flashcards must be directly related to the content above.
    3. Cover different topics and concepts from the text.
    4. Focus on key concepts, facts, and important details.

    RESPONSE FORMAT (JSON only):
    {{
        "flashcards": [
            {{
                "front": "What is [key concept]?",
                "back": "[Definition/Explanation of key concept]"
            }}
        ]
    }}

    Generate {num_flashcards} flashcards now. Return ONLY the JSON response, no additional text.
    """

    payload = {
        "model": "llama3-8b-8192",
        "messages": [
            {
                "role": "system",
                "content": "You are an expert educational content creator specializing in generating high-quality flashcards from text content. Always respond with valid JSON format only."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        "temperature": 0.3,
        "max_tokens": 2500
    }

    try:
        print(f"Calling GROQ API to generate {num_flashcards} flashcards...")
        response = requests.post(GROQ_API_URL, headers=headers, json=payload)

        if response.status_code == 200:
            result = response.json()
            content = result['choices'][0]['message']['content'].strip()

            if content.startswith('```json'):
                content = content.replace('```json', '').replace('```', '').strip()

            try:
                flashcard_data = json.loads(content)
                if 'flashcards' in flashcard_data and len(flashcard_data['flashcards']) > 0:
                    print(f"Successfully generated {len(flashcard_data['flashcards'])} flashcards")
                    return flashcard_data['flashcards']
                else:
                    print("No flashcards found in API response")
                    return create_content_based_flashcards(text, num_flashcards)
            except json.JSONDecodeError as e:
                print(f"JSON parsing error for flashcards: {e}")
                return create_content_based_flashcards(text, num_flashcards)
        else:
            print(f"GROQ API error for flashcards: {response.status_code} - {response.text}")
            return create_content_based_flashcards(text, num_flashcards)

    except Exception as e:
        print(f"Error calling GROQ API for flashcards: {e}")
        return create_content_based_flashcards(text, num_flashcards)

def create_content_based_flashcards(text, num_flashcards=10):
    """Fallback: Create simple flashcards based on text content"""
    flashcards = []
    sentences = [s.strip() for s in text.split('.') if s.strip()]

    for i in range(min(num_flashcards, len(sentences) // 2)):
        front = sentences[i * 2]
        back = sentences[i * 2 + 1] if (i * 2 + 1) < len(sentences) else "No further content."
        flashcards.append({"front": front, "back": back})

    if not flashcards and sentences:
        # If not enough pairs, create single flashcard from first sentence
        flashcards.append({"front": sentences[0], "back": "Key point from the text."})

    return flashcards[:num_flashcards]

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['GET', 'POST'])
def upload():
    if request.method == 'POST':
        # Check if file is present
        if 'file' not in request.files:
            flash('No file selected', 'error')
            return redirect(request.url)
        
        file = request.files['file']
        
        if file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            
            # Extract text based on file type
            file_ext = filename.rsplit('.', 1)[1].lower()
            text = ""
            print(f"Processing file: {filename} (type: {file_ext})")
            file_size = os.path.getsize(file_path)
            print(f"File size: {file_size/1024:.2f} KB")
            
            # Validate file content before processing
            is_valid, validation_error = validate_file_content(file_path, file_ext)
            if not is_valid:
                flash(validation_error, 'error')
                os.remove(file_path)
                return redirect(request.url)
            
            try:
                if file_ext == 'pdf':
                    # First try to extract text normally
                    text = extract_text_from_pdf(file_path)
                    
                    # If minimal text was extracted, try full page OCR
                    if not text.strip() or len(text.strip()) < 50:
                        print("Minimal text found in PDF, attempting full page OCR...")
                        if setup_tesseract():
                            try:
                                images = convert_from_path(file_path, dpi=300)
                                ocr_text = ""
                                
                                for i, image in enumerate(images):
                                    processed_img = preprocess_image_for_ocr(image)
                                    page_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                                    if page_text.strip():
                                        ocr_text += f"\n[Page {i + 1}]:\n{page_text}\n"
                                        print(f"OCR Page {i + 1}: Extracted {len(page_text)} characters")
                                
                                if ocr_text.strip():
                                    text = ocr_text
                                    print("Successfully extracted text using full page OCR")
                            except Exception as e:
                                print(f"Full page OCR failed: {e}")
                
                elif file_ext in ['docx', 'doc']:
                    text = extract_text_from_docx(file_path)
                    # If minimal text was found, try OCR on any images in the document
                    if not text.strip() or len(text.strip()) < 50:
                        if setup_tesseract() and file_ext == 'docx':
                            image_text = extract_images_from_docx(file_path)
                            if image_text:
                                text = image_text
                
                elif file_ext == 'pptx':
                    text = extract_text_from_pptx(file_path)
                    # If minimal text was found, try OCR on any images in the presentation
                    if not text.strip() or len(text.strip()) < 50:
                        if setup_tesseract():
                            image_text = extract_images_from_pptx(file_path)
                            if image_text:
                                text = image_text
                
                elif file_ext in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
                    if setup_tesseract():
                        text = extract_text_from_image(file_path)
                    else:
                        flash('Tesseract OCR is not installed. Please install it to process images.', 'error')
                        os.remove(file_path)
                        return redirect(request.url)
                
                print(f"Total extracted text length: {len(text)} characters")
                
                if not text.strip():
                    error_msg = 'No text could be extracted from the file. This might be because:\n'
                    error_msg += '- The file is password protected\n'
                    error_msg += '- The file contains scanned images without OCR\n'
                    error_msg += '- The file is corrupted or empty\n'
                    error_msg += '- The file contains only non-text content (like images)\n\n'
                    error_msg += 'Please ensure the file contains readable text and try again.'
                    flash(error_msg, 'error')
                    os.remove(file_path)
                    return redirect(request.url)
                
                # Generate MCQs
                mcqs = generate_mcqs_with_groq(text, num_questions=10)
                
                # Generate Flashcards
                flashcards = generate_flashcards_with_groq(text, num_flashcards=10)
                
                # Store generated content in session for temporary access
                # This overwrites any previous session data with new content
                session['mcqs'] = mcqs
                session['flashcards'] = flashcards
                session['filename'] = filename
                
                print(f"Stored {len(mcqs)} MCQs and {len(flashcards)} flashcards in session")
                
                # Clean up uploaded file
                os.remove(file_path)
                
                return render_template('quizmc.html', mcqs=mcqs, flashcards=flashcards, filename=filename)
                
            except (PermissionError, OSError) as e:
                print(f"File access error: {e}")
                flash('Unable to access the file. The file might be in use or you may not have permission to read it.', 'error')
                if os.path.exists(file_path):
                    os.remove(file_path)
                return redirect(request.url)
            except Exception as e:
                print(f"Exception during processing: {e}")
                error_msg = 'Error processing file. Please check that:\n'
                error_msg += '- The file is not corrupted\n'
                error_msg += '- The file format matches its extension\n'
                error_msg += '- The file is not encrypted or password protected'
                flash(error_msg, 'error')
                if os.path.exists(file_path):
                    os.remove(file_path)
                return redirect(request.url)
        
        else:
            flash('Invalid file type. Please upload PDF, DOC, DOCX, PPTX, or image files.', 'error')
    
    return render_template('upload.html')

@app.route('/check_answer', methods=['POST'])
def check_answer():
    """API endpoint to check if selected answer is correct"""
    data = request.json
    question_index = data.get('question_index')
    selected_option = data.get('selected_option')
    correct_answer = data.get('correct_answer')
    
    is_correct = selected_option == correct_answer
    
    return jsonify({
        'is_correct': is_correct,
        'correct_answer': correct_answer
    })



@app.route('/session_data')
def get_session_data():
    """API endpoint to retrieve current session's MCQs and flashcards"""
    try:
        # Get data from session storage
        mcqs = session.get('mcqs', [])
        flashcards = session.get('flashcards', [])
        filename = session.get('filename', 'No file')
        
        return jsonify({
            'success': True,
            'mcqs': mcqs,
            'flashcards': flashcards,
            'filename': filename,
            'mcq_count': len(mcqs),
            'flashcard_count': len(flashcards)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/clear_session')
def clear_session_data():
    """Clear all session data"""
    session.pop('mcqs', None)
    session.pop('flashcards', None)
    session.pop('filename', None)
    
    return jsonify({
        'success': True,
        'message': 'Session data cleared'
    })

# 🧠 Main chat with session data (unchanged)
@app.route('/quizpasa_chat', methods=['POST'])
def quizpasa_chat():
    """Handle chat messages and provide responses with session context"""
    try:
        data = request.get_json()
        user_message = data.get('message', '')
        conversation_history = data.get('conversation_history', [])
        
        # Get current session data
        session_data = {
            'mcqs': session.get('mcqs', []),
            'flashcards': session.get('flashcards', [])
        }
        
        # Get response from QuizPasa with session context
        response = handle_chat_message(user_message, conversation_history, session_data)
        return jsonify(response)
        
    except Exception as e:
        return jsonify({
            'success': False,
            'message': 'An error occurred while processing your message.',
            'error': str(e)
        }), 500


# 🤖 Simpler fallback chat route (renamed)
@app.route('/quizpasa_chat_simple', methods=['POST'])
def quizpasa_chat_simple():
    """Handle basic chatbot messages"""
    try:
        data = request.get_json()
        user_message = data.get('message', '').strip()
        conversation_history = data.get('history', [])
        
        if not user_message:
            return jsonify({
                'success': False,
                'error': 'No message provided'
            }), 400
        
        # Handle special commands
        if user_message.lower() in ['hi', 'hello', 'hey']:
            response = get_welcome_message()
        else:
            # Get response from QuizPasa
            response = handle_chat_message(user_message, conversation_history)
        
        return jsonify(response)
        
    except Exception as e:
        return jsonify({
            'success': False,
            'message': 'Sorry, I encountered an error. Please try again!',
            'error': str(e)
        }), 500



# if __name__ == '__main__':
#     port = int(os.environ.get('PORT', 5000))
#     app.run(host='0.0.0.0', port=port)

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
